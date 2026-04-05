#!/usr/bin/env python3
"""Generate a beautiful 15-slide FinanceFlow PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # neon cyan
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # emerald
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)   # soft purple
ACCENT_ORANGE= RGBColor(0xFB, 0xBF, 0x24)   # amber
ACCENT_PINK  = RGBColor(0xF4, 0x72, 0xB6)   # pink
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)   # slate-800

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ─────────────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    """Fill the entire slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color, corner_radius=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def text_box(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def title_subtitle(slide, title, subtitle, title_size=44, sub_size=20,
                   title_color=WHITE, sub_color=LIGHT_GRAY):
    """Add a title + subtitle to a slide."""
    text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             title, title_size, title_color, bold=True, font_name="Calibri Light")
    accent_line(slide, Inches(0.8), Inches(1.45), Inches(2), ACCENT_BLUE)
    text_box(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.7),
             subtitle, sub_size, sub_color, font_name="Calibri")


def bullet_card(slide, left, top, width, height, title, bullets,
                card_color=CARD_BG, title_color=ACCENT_BLUE, bullet_color=LIGHT_GRAY):
    """Card with title and bullet points."""
    card = rect(slide, left, top, width, height, card_color, 0.05)
    # title
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25),
                                      width - Inches(0.6), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"
    # bullets
    txBox2 = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.85),
                                       width - Inches(0.6), height - Inches(1.1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = f"  {b}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = bullet_color
        p2.font.name = "Calibri"
        p2.space_after = Pt(6)
    return card


def icon_card(slide, left, top, width, height, icon_text, label,
              value, card_color=CARD_BG, icon_color=ACCENT_BLUE):
    """Stat card with icon, label, and value."""
    card = rect(slide, left, top, width, height, card_color, 0.08)
    # icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.3), top + Inches(0.3),
        Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_color
    circle.line.fill.background()
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.35),
                                      Inches(0.6), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # value
    text_box(slide, left + Inches(1.1), top + Inches(0.25), width - Inches(1.4), Inches(0.45),
             value, 22, WHITE, bold=True)
    # label
    text_box(slide, left + Inches(1.1), top + Inches(0.65), width - Inches(1.4), Inches(0.35),
             label, 13, MID_GRAY)
    return card


def slide_number(slide, num):
    """Add a slide number."""
    text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
             str(num), 11, MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s)
# Large gradient-style accent bar at top
accent_line(s, Inches(0), Inches(0), W, ACCENT_BLUE)
# Logo area
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(1.5), Inches(1.7), Inches(1.7))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT_BLUE
circle.line.fill.background()
txBox = s.shapes.add_textbox(Inches(5.8), Inches(1.7), Inches(1.7), Inches(1.3))
p = txBox.text_frame.paragraphs[0]
p.text = "FF"
p.font.size = Pt(52)
p.font.color.rgb = DARK_BG
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

text_box(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2),
         "FinanceFlow", 56, WHITE, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")
text_box(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.8),
         "Personal Finance Management Application", 24, ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)
accent_line(s, Inches(5.5), Inches(5.5), Inches(2.3), ACCENT_BLUE)
text_box(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
         "A Modern Full-Stack Solution for Smart Money Management", 16, LIGHT_GRAY,
         alignment=PP_ALIGN.CENTER)
text_box(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
         "Web Programming Project  |  2025", 14, MID_GRAY,
         alignment=PP_ALIGN.CENTER)
slide_number(s, 1)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "The Problem", "Why personal finance management matters")

problems = [
    ("78%", "of people live paycheck to paycheck", ACCENT_PINK),
    ("65%", "don't track their monthly spending", ACCENT_ORANGE),
    ("43%", "have no emergency savings fund", ACCENT_PURPLE),
]
for i, (stat, desc, color) in enumerate(problems):
    left = Inches(0.8 + i * 4.1)
    card = rect(s, left, Inches(2.8), Inches(3.7), Inches(2.2), CARD_BG, 0.06)
    text_box(s, left + Inches(0.3), Inches(3.0), Inches(3.1), Inches(1),
             stat, 48, color, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.3), Inches(4.0), Inches(3.1), Inches(0.8),
             desc, 16, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

text_box(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1),
         "People need an intuitive, all-in-one tool to track income, manage expenses,\nset savings goals, and visualize their financial health — without complexity.",
         16, MID_GRAY, alignment=PP_ALIGN.CENTER)
slide_number(s, 2)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — Our Solution
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Our Solution", "FinanceFlow — Your complete financial companion")

features = [
    ("Track", "Income & Expenses", "Real-time tracking with\ncategory breakdowns", ACCENT_BLUE),
    ("Analyze", "Reports & Analytics", "Visual charts, trends,\nand financial health score", ACCENT_GREEN),
    ("Plan", "Savings Goals", "Set targets and track\nprogress automatically", ACCENT_PURPLE),
    ("Secure", "Authentication", "Protected accounts with\nsession management", ACCENT_ORANGE),
]
for i, (icon, title, desc, color) in enumerate(features):
    left = Inches(0.5 + i * 3.15)
    card = rect(s, left, Inches(2.7), Inches(2.9), Inches(3.5), CARD_BG, 0.06)
    # colored top bar
    accent_line(s, left, Inches(2.7), Inches(2.9), color)
    text_box(s, left + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.5),
             icon, 28, color, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.2), Inches(3.5), Inches(2.5), Inches(0.5),
             title, 17, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.2), Inches(4.2), Inches(2.5), Inches(1),
             desc, 13, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(s, 3)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — Tech Stack
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Tech Stack", "Modern, lightweight, and production-ready")

# Frontend card
bullet_card(s, Inches(0.5), Inches(2.5), Inches(3.8), Inches(4.2),
            "Frontend", [
                "HTML5 / CSS3 / Vanilla JS",
                "Tailwind CSS (utility-first)",
                "Chart.js (data visualization)",
                "Material Icons (Google)",
                "SVG custom charts",
                "LocalStorage for caching",
            ], title_color=ACCENT_BLUE)

# Backend card
bullet_card(s, Inches(4.7), Inches(2.5), Inches(3.8), Inches(4.2),
            "Backend", [
                "Node.js v18+",
                "Express.js v5.2",
                "Helmet (security headers)",
                "CORS & Rate Limiting",
                "ES6+ Modules (import/export)",
                "dotenv configuration",
            ], title_color=ACCENT_GREEN)

# Database card
bullet_card(s, Inches(8.9), Inches(2.5), Inches(3.8), Inches(4.2),
            "Database", [
                "MySQL 8.0+ (primary)",
                "In-Memory fallback (dev)",
                "Connection pooling (x10)",
                "Auto-detection on startup",
                "Migration scripts included",
                "Foreign key constraints",
            ], title_color=ACCENT_PURPLE)

slide_number(s, 4)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — System Architecture
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "System Architecture", "Hybrid storage with intelligent fallback")

# Architecture diagram using shapes
layers = [
    ("Browser (HTML/JS/Tailwind)", Inches(4), Inches(2.5), Inches(5.3), Inches(0.9), ACCENT_BLUE),
    ("Express.js API Server", Inches(4), Inches(3.8), Inches(5.3), Inches(0.9), ACCENT_GREEN),
    ("Data Access Layer (Abstraction)", Inches(4), Inches(5.1), Inches(5.3), Inches(0.9), ACCENT_ORANGE),
]
for label, l, t, w, h, color in layers:
    card = rect(s, l, t, w, h, CARD_BG, 0.08)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    text_box(s, l, t + Inches(0.15), w, Inches(0.6),
             label, 18, color, bold=True, alignment=PP_ALIGN.CENTER)

# Arrows between layers (simple rectangles)
for y in [Inches(3.4), Inches(4.7)]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(6.5), y, Pt(3), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MID_GRAY
    arrow.line.fill.background()

# DB options
db1 = rect(s, Inches(2.5), Inches(6.3), Inches(3.5), Inches(0.8), CARD_BG, 0.08)
db1.line.color.rgb = ACCENT_PURPLE
db1.line.width = Pt(2)
text_box(s, Inches(2.5), Inches(6.4), Inches(3.5), Inches(0.6),
         "MySQL 8.0+", 16, ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

db2 = rect(s, Inches(7.3), Inches(6.3), Inches(3.5), Inches(0.8), CARD_BG, 0.08)
db2.line.color.rgb = ACCENT_PINK
db2.line.width = Pt(2)
text_box(s, Inches(7.3), Inches(6.4), Inches(3.5), Inches(0.6),
         "In-Memory Arrays", 16, ACCENT_PINK, bold=True, alignment=PP_ALIGN.CENTER)

# Side labels
text_box(s, Inches(0.5), Inches(2.6), Inches(3.2), Inches(0.8),
         "Client-side rendering\nLocalStorage caching\nResponsive UI", 13, LIGHT_GRAY)
text_box(s, Inches(0.5), Inches(3.9), Inches(3.2), Inches(0.8),
         "RESTful API endpoints\nJSON responses\nSession auth", 13, LIGHT_GRAY)
text_box(s, Inches(0.5), Inches(5.2), Inches(3.2), Inches(0.8),
         "Auto-detects MySQL\nSeamless fallback\nUnified interface", 13, LIGHT_GRAY)

slide_number(s, 5)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — Dashboard
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Dashboard", "Real-time financial overview at a glance")

# Stat cards
stats = [
    ("$", "Total Income", "$24,500", ACCENT_GREEN),
    ("-", "Total Expenses", "$18,200", ACCENT_PINK),
    ("=", "Net Balance", "$6,300", ACCENT_BLUE),
    ("%", "Savings Rate", "25.7%", ACCENT_PURPLE),
]
for i, (icon, label, val, color) in enumerate(stats):
    icon_card(s, Inches(0.5 + i * 3.15), Inches(2.5), Inches(2.9), Inches(1.2),
              icon, label, val, icon_color=color)

# Feature cards
features = [
    ("Portfolio Chart", "Interactive SVG line chart\nwith 7/30/90/180-day views\nand dynamic data scaling", ACCENT_BLUE),
    ("Financial Health Score", "Intelligent 0-100 score based on:\n- Savings rate (40% weight)\n- Expense consistency (30%)\n- Income growth (30%)", ACCENT_GREEN),
    ("Recent Transactions", "Live feed of latest income\nand expense entries with\ncategory tags and amounts", ACCENT_ORANGE),
    ("Savings Goals", "Visual progress bars for\nmultiple concurrent goals\nwith target tracking", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(features):
    left = Inches(0.5 + i * 3.15)
    card = rect(s, left, Inches(4.2), Inches(2.9), Inches(2.8), CARD_BG, 0.06)
    accent_line(s, left, Inches(4.2), Inches(2.9), color)
    text_box(s, left + Inches(0.2), Inches(4.45), Inches(2.5), Inches(0.4),
             title, 16, color, bold=True)
    text_box(s, left + Inches(0.2), Inches(4.95), Inches(2.5), Inches(1.8),
             desc, 12, LIGHT_GRAY)

slide_number(s, 6)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — Expense Management
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Expense Management", "Track, categorize, and analyze every rupee")

# Categories with colors
cats = [
    ("Housing", "Rent, utilities, maintenance", ACCENT_BLUE, "40%"),
    ("Food", "Groceries, dining out", ACCENT_GREEN, "22%"),
    ("Transport", "Fuel, public transit, rides", ACCENT_ORANGE, "15%"),
    ("Entertainment", "Movies, subscriptions", ACCENT_PINK, "10%"),
    ("Technology", "Gadgets, software, cloud", ACCENT_PURPLE, "8%"),
    ("Other", "Miscellaneous spending", MID_GRAY, "5%"),
]
for i, (cat, desc, color, pct) in enumerate(cats):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.15)
    top = Inches(2.5 + row * 1.6)
    card = rect(s, left, top, Inches(3.8), Inches(1.3), CARD_BG, 0.06)
    # color indicator
    indicator = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left + Inches(0.15), top + Inches(0.3),
                                    Pt(6), Inches(0.7))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    text_box(s, left + Inches(0.4), top + Inches(0.2), Inches(2.2), Inches(0.35),
             cat, 17, WHITE, bold=True)
    text_box(s, left + Inches(0.4), top + Inches(0.55), Inches(2.2), Inches(0.35),
             desc, 12, MID_GRAY)
    text_box(s, left + Inches(2.8), top + Inches(0.3), Inches(0.8), Inches(0.5),
             pct, 24, color, bold=True, alignment=PP_ALIGN.RIGHT)

# Key features
text_box(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.8),
         "Donut chart visualization  |  Month-over-month trends  |  Smart spending suggestions  |  Search & filter",
         14, ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
slide_number(s, 7)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — Income Tracking
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Income Tracking", "Monitor all revenue streams in one place")

# Income categories
inc_cats = [
    ("Salary", "Primary employment income", ACCENT_GREEN, "$15,000"),
    ("Freelance", "Side projects & contracts", ACCENT_BLUE, "$5,200"),
    ("Investment", "Dividends & returns", ACCENT_PURPLE, "$3,100"),
    ("Other", "Miscellaneous income", ACCENT_ORANGE, "$1,200"),
]
for i, (cat, desc, color, amt) in enumerate(inc_cats):
    left = Inches(0.5 + i * 3.15)
    card = rect(s, left, Inches(2.5), Inches(2.9), Inches(1.6), CARD_BG, 0.06)
    accent_line(s, left, Inches(2.5), Inches(2.9), color)
    text_box(s, left + Inches(0.2), Inches(2.75), Inches(2.5), Inches(0.4),
             cat, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.2), Inches(3.15), Inches(2.5), Inches(0.4),
             amt, 28, color, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.2), Inches(3.55), Inches(2.5), Inches(0.4),
             desc, 12, MID_GRAY, alignment=PP_ALIGN.CENTER)

# Features
bullet_card(s, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.5),
            "Key Features", [
                "Add, view, filter, and manage income entries",
                "Bar chart visualization for last 6 months",
                "Month-over-month comparison analysis",
                "Search by description and category filtering",
                "Real-time totals and trend indicators",
            ], title_color=ACCENT_GREEN)

bullet_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
            "Insights Generated", [
                "Monthly income trend analysis",
                "Category distribution breakdown",
                "Income growth rate calculation",
                "Highest earning month identification",
                "Projected annual income estimation",
            ], title_color=ACCENT_BLUE)

slide_number(s, 8)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — Analytics & Reports
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Analytics & Reports", "Data-driven insights powered by Chart.js")

charts = [
    ("Line Charts", "Weekly & monthly\nspending trends over time", ACCENT_BLUE),
    ("Bar Charts", "Income vs. expenses\nside-by-side comparison", ACCENT_GREEN),
    ("Pie Charts", "Category distribution\npercentage breakdown", ACCENT_PURPLE),
    ("Health Score", "Combined metric from\nsavings, expenses & growth", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(charts):
    left = Inches(0.5 + i * 3.15)
    card = rect(s, left, Inches(2.5), Inches(2.9), Inches(2.2), CARD_BG, 0.06)
    # Chart icon placeholder
    icon_rect = rect(s, left + Inches(0.8), Inches(2.7), Inches(1.3), Inches(0.8), color, 0.1)
    text_box(s, left + Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.4),
             title, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    text_box(s, left + Inches(0.2), Inches(4.1), Inches(2.5), Inches(0.7),
             desc, 13, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Analytics insights
insights = [
    "Real-time data processing from live transactions",
    "Custom date range filtering (7, 30, 90, 180 days)",
    "Financial health algorithm combining 3 weighted metrics",
    "Interactive tooltips and responsive chart rendering",
]
card = rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), CARD_BG, 0.04)
text_box(s, Inches(0.8), Inches(5.3), Inches(3), Inches(0.4),
         "Analytics Highlights", 18, ACCENT_BLUE, bold=True)
for i, insight in enumerate(insights):
    col = i % 2
    row = i // 2
    text_box(s, Inches(0.8 + col * 6), Inches(5.8 + row * 0.5),
             Inches(5.8), Inches(0.4),
             f"  {insight}", 13, LIGHT_GRAY)

slide_number(s, 9)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — Savings Goals
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Savings Goals", "Set targets and watch your progress grow")

goals = [
    ("New Porsche 911", "$195,000", "$126,750", "65%", ACCENT_BLUE),
    ("Tokyo Trip", "$5,000", "$4,600", "92%", ACCENT_GREEN),
    ("Emergency Fund", "$10,000", "$7,500", "75%", ACCENT_PURPLE),
]
for i, (name, target, saved, pct, color) in enumerate(goals):
    top = Inches(2.5 + i * 1.6)
    card = rect(s, Inches(1), top, Inches(11.3), Inches(1.35), CARD_BG, 0.04)
    text_box(s, Inches(1.3), top + Inches(0.15), Inches(3), Inches(0.35),
             name, 20, WHITE, bold=True)
    text_box(s, Inches(1.3), top + Inches(0.5), Inches(3), Inches(0.3),
             f"Target: {target}  |  Saved: {saved}", 13, MID_GRAY)
    # Progress bar background
    bar_bg = rect(s, Inches(5.5), top + Inches(0.45), Inches(5.5), Inches(0.35), RGBColor(0x33, 0x3D, 0x50), 0.5)
    # Progress bar fill
    pct_val = int(pct.replace('%', '')) / 100
    bar_fill = rect(s, Inches(5.5), top + Inches(0.45), Inches(5.5 * pct_val), Inches(0.35), color, 0.5)
    # Percentage
    text_box(s, Inches(11.2), top + Inches(0.15), Inches(1), Inches(0.4),
             pct, 22, color, bold=True, alignment=PP_ALIGN.RIGHT)

# Note
text_box(s, Inches(1), Inches(7.3), Inches(11.3), Inches(0.5),
         "Goals update automatically as you log income and expenses  |  Multiple concurrent goals supported",
         13, MID_GRAY, alignment=PP_ALIGN.CENTER)
slide_number(s, 10)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 11 — Authentication & Security
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Authentication & Security", "Protected accounts with modern security practices")

# Left side - Auth features
bullet_card(s, Inches(0.5), Inches(2.5), Inches(5.8), Inches(4.5),
            "Authentication System", [
                "Email & password registration with validation",
                "Token-based session management",
                "Secure localStorage token storage",
                "Duplicate email detection on signup",
                "Demo account for instant testing",
                "Glassmorphic login UI design",
                "Auto-redirect on expired sessions",
            ], title_color=ACCENT_BLUE)

# Right side - Security
bullet_card(s, Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.5),
            "Security Features", [
                "Helmet.js security headers",
                "CORS configuration",
                "API rate limiting protection",
                "Input validation on all endpoints",
                "SQL injection prevention",
                "Two-factor authentication UI (ready)",
                "Environment-based configuration",
            ], title_color=ACCENT_GREEN)

slide_number(s, 11)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 12 — Settings & Personalization
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "Settings & Personalization", "Customize your FinanceFlow experience")

settings = [
    ("Currency", "INR, USD, EUR, GBP, JPY\nSwitch anytime, app-wide", ACCENT_BLUE),
    ("Dark Mode", "Toggle dark/light theme\nPersists across sessions", ACCENT_PURPLE),
    ("Budget Goal", "Set monthly budget target\nTrack progress automatically", ACCENT_GREEN),
    ("Profile", "Avatar upload, name & email\nMembership since date", ACCENT_ORANGE),
    ("Notifications", "Email alerts toggle\nStay informed on activity", ACCENT_PINK),
    ("2FA Ready", "Two-factor auth UI built\nReady for backend integration", MID_GRAY),
]
for i, (title, desc, color) in enumerate(settings):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.15)
    top = Inches(2.5 + row * 2.3)
    card = rect(s, left, top, Inches(3.8), Inches(2.0), CARD_BG, 0.06)
    # color dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.3),
                              Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    text_box(s, left + Inches(0.75), top + Inches(0.25), Inches(2.8), Inches(0.35),
             title, 18, WHITE, bold=True)
    text_box(s, left + Inches(0.75), top + Inches(0.65), Inches(2.8), Inches(1),
             desc, 13, LIGHT_GRAY)

slide_number(s, 12)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 13 — API Endpoints
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "RESTful API Design", "Clean, consistent, and well-structured endpoints")

endpoints = [
    ("GET", "/api/dashboard", "Dashboard stats, transactions, goals", ACCENT_GREEN),
    ("GET", "/api/transactions", "Filtered transaction list", ACCENT_GREEN),
    ("POST", "/api/transactions", "Create new transaction", ACCENT_BLUE),
    ("DELETE", "/api/transactions/:id", "Remove a transaction", ACCENT_PINK),
    ("GET", "/api/analytics", "Charts data, trends, insights", ACCENT_GREEN),
    ("GET", "/api/goals", "User savings goals", ACCENT_GREEN),
    ("GET", "/api/settings", "User preferences", ACCENT_GREEN),
    ("PUT", "/api/settings", "Update user settings", ACCENT_ORANGE),
    ("POST", "/api/auth/login", "User authentication", ACCENT_BLUE),
    ("POST", "/api/auth/signup", "User registration", ACCENT_BLUE),
]

# Table header
card = rect(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.6), ACCENT_BLUE, 0.03)
text_box(s, Inches(1.0), Inches(2.35), Inches(1.2), Inches(0.4), "Method", 14, DARK_BG, bold=True)
text_box(s, Inches(2.5), Inches(2.35), Inches(3.5), Inches(0.4), "Endpoint", 14, DARK_BG, bold=True)
text_box(s, Inches(6.5), Inches(2.35), Inches(5.5), Inches(0.4), "Description", 14, DARK_BG, bold=True)

for i, (method, path, desc, color) in enumerate(endpoints):
    top = Inches(3.0 + i * 0.42)
    bg_color = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x33)
    row = rect(s, Inches(0.8), top, Inches(11.7), Inches(0.4), bg_color, 0.02)
    text_box(s, Inches(1.0), top + Inches(0.03), Inches(1.2), Inches(0.35),
             method, 12, color, bold=True)
    text_box(s, Inches(2.5), top + Inches(0.03), Inches(3.5), Inches(0.35),
             path, 12, WHITE, font_name="Consolas")
    text_box(s, Inches(6.5), top + Inches(0.03), Inches(5.5), Inches(0.35),
             desc, 12, LIGHT_GRAY)

slide_number(s, 13)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 14 — Unique Selling Points
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_subtitle(s, "What Makes Us Different", "Key differentiators of FinanceFlow")

usps = [
    ("Zero-Config Deployment",
     "Works immediately with in-memory storage. No database setup needed. Just run and go.",
     ACCENT_GREEN),
    ("Hybrid Storage Engine",
     "Auto-detects MySQL availability. Graceful fallback to in-memory. One unified API layer.",
     ACCENT_BLUE),
    ("Premium Glassmorphic UI",
     "Neon glow effects, dark mode, responsive design. A premium look without heavy frameworks.",
     ACCENT_PURPLE),
    ("Financial Health AI",
     "Intelligent 0-100 score combining savings rate, expense patterns, and income growth.",
     ACCENT_ORANGE),
    ("Lightweight & Fast",
     "No frontend framework overhead. Vanilla JS + Tailwind = blazing fast load times.",
     ACCENT_PINK),
    ("Production Ready",
     "Helmet security, rate limiting, CORS, connection pooling. Ready for real users.",
     ACCENT_BLUE),
]
for i, (title, desc, color) in enumerate(usps):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.15)
    top = Inches(2.5 + row * 2.3)
    card = rect(s, left, top, Inches(3.8), Inches(2.0), CARD_BG, 0.06)
    accent_line(s, left, top, Inches(3.8), color)
    text_box(s, left + Inches(0.25), top + Inches(0.25), Inches(3.3), Inches(0.35),
             title, 17, color, bold=True)
    text_box(s, left + Inches(0.25), top + Inches(0.7), Inches(3.3), Inches(1.1),
             desc, 13, LIGHT_GRAY)

slide_number(s, 14)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 15 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
accent_line(s, Inches(0), Inches(0), W, ACCENT_BLUE)

# Decorative circles
for cx, cy, sz, clr in [
    (Inches(1.5), Inches(1.5), Inches(2), ACCENT_BLUE),
    (Inches(10), Inches(5.5), Inches(1.5), ACCENT_PURPLE),
    (Inches(0.3), Inches(5), Inches(1), ACCENT_GREEN),
    (Inches(11.5), Inches(1.2), Inches(0.8), ACCENT_PINK),
]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.fill.fore_color.brightness = 0.7
    c.line.fill.background()

text_box(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
         "Thank You!", 60, WHITE, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

accent_line(s, Inches(5.5), Inches(3.8), Inches(2.3), ACCENT_BLUE)

text_box(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
         "FinanceFlow — Smart Money Management Made Simple", 22, ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)

text_box(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
         "Questions & Discussion", 20, LIGHT_GRAY,
         alignment=PP_ALIGN.CENTER)

# Demo info card
card = rect(s, Inches(4), Inches(5.9), Inches(5.3), Inches(1), CARD_BG, 0.06)
text_box(s, Inches(4.3), Inches(6.0), Inches(4.7), Inches(0.4),
         "Demo: demo@financeflow.local / demo123", 14, MID_GRAY,
         alignment=PP_ALIGN.CENTER)
text_box(s, Inches(4.3), Inches(6.4), Inches(4.7), Inches(0.4),
         "Run:  node server/src/full-server.mjs", 14, ACCENT_GREEN,
         alignment=PP_ALIGN.CENTER, font_name="Consolas")

slide_number(s, 15)


# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FinanceFlow_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
